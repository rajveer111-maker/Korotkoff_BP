import os
import subprocess
import sys

# Ensure python-pptx is installed
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
except ImportError:
    print("Installing python-pptx...")
    subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx"])
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN

def create_presentation():
    prs = Presentation()
    
    # Slide Layouts
    TITLE_SLIDE = 0
    BULLET_SLIDE = 1
    PIC_SLIDE = 5 # Title only, good for full pictures
    
    # ---------------------------------------------------------
    # Slide 1: Title
    # ---------------------------------------------------------
    slide = prs.slides.add_slide(prs.slide_layouts[TITLE_SLIDE])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "High-Fidelity Radar Radiomyography (RMG)"
    subtitle.text = "Non-Invasive Korotkoff Sound Detection for Blood Pressure Measurement\n\nAdvanced Multi-Domain Validation Pipeline"

    # ---------------------------------------------------------
    # Slide 2: Introduction & Motivation
    # ---------------------------------------------------------
    slide = prs.slides.add_slide(prs.slide_layouts[BULLET_SLIDE])
    slide.shapes.title.text = "Introduction & Motivation"
    tf = slide.placeholders[1].text_frame
    tf.text = "Limitations of traditional blood pressure (BP) measurement:"
    p = tf.add_paragraph(); p.text = "Stethoscopes require trained human hearing and are prone to artifact noise."; p.level = 1
    p = tf.add_paragraph(); p.text = "Oscillometric cuffs use generic algorithms that guess SYS/DIA points."; p.level = 1
    p = tf.add_paragraph(); p.text = "Our Solution: RF Radar Sensing (Radiomyography)"; p.level = 0
    p = tf.add_paragraph(); p.text = "Use a USRP software-defined radio to measure sub-millimeter arterial vibrations."; p.level = 1
    p = tf.add_paragraph(); p.text = "Mathematically isolate true mechanical Korotkoff 'snaps' during deflation."; p.level = 1

    # ---------------------------------------------------------
    # Slide 3: Hardware & System Architecture
    # ---------------------------------------------------------
    slide = prs.slides.add_slide(prs.slide_layouts[BULLET_SLIDE])
    slide.shapes.title.text = "Hardware & System Architecture"
    tf = slide.placeholders[1].text_frame
    tf.text = "Software-Defined Radio Setup:"
    p = tf.add_paragraph(); p.text = "Device: USRP B210 (Dual Channel capability, strictly isolated to Port A)."; p.level = 1
    p = tf.add_paragraph(); p.text = "Frequency: 0.9 GHz (Optimized for tissue penetration)."; p.level = 1
    p = tf.add_paragraph(); p.text = "Sampling Rate: 10 kHz (Downsampled from 1 MHz for high-resolution audio-band tracking)."; p.level = 1
    p = tf.add_paragraph(); p.text = "Gain Settings: TX Gain = 40 dB, RX Gain = 45-60 dB (Tuned to prevent saturation while capturing micro-vibrations)."; p.level = 1

    # ---------------------------------------------------------
    # Slide 4: Data Acquisition & Protocol
    # ---------------------------------------------------------
    slide = prs.slides.add_slide(prs.slide_layouts[BULLET_SLIDE])
    slide.shapes.title.text = "Experimental Protocol"
    tf = slide.placeholders[1].text_frame
    tf.text = "Cuff Deflation Profile:"
    p = tf.add_paragraph(); p.text = "The subject wears a standard BP cuff over the brachial artery."; p.level = 1
    p = tf.add_paragraph(); p.text = "Phase 1 (Occlusion): Cuff inflated above Systolic pressure (>140 mmHg). Blood flow stops."; p.level = 1
    p = tf.add_paragraph(); p.text = "Phase 2 (Korotkoff Window): Cuff slowly deflates. Blood spurts through the compressed artery, causing mechanical vibrations (10-50 Hz)."; p.level = 1
    p = tf.add_paragraph(); p.text = "Phase 3 (Free Flow): Cuff drops below Diastolic pressure. Smooth flow resumes, vibrations cease."; p.level = 1

    # ---------------------------------------------------------
    # Slide 5: Signal Pre-Processing & Physical Units
    # ---------------------------------------------------------
    slide = prs.slides.add_slide(prs.slide_layouts[BULLET_SLIDE])
    slide.shapes.title.text = "Pre-Processing & Phase Detrending"
    tf = slide.placeholders[1].text_frame
    tf.text = "Extracting physical movement from raw RF waves:"
    p = tf.add_paragraph(); p.text = "I/Q Processing: Raw In-Phase and Quadrature data are AC-coupled to remove DC offsets."; p.level = 1
    p = tf.add_paragraph(); p.text = "Phase Unwrapping: Calculates absolute angle, but suffers from infinite accumulation drift."; p.level = 1
    p = tf.add_paragraph(); p.text = "Linear Detrending: Crucial algorithm applied to flatten the baseline phase drift."; p.level = 1
    p = tf.add_paragraph(); p.text = "Physical Conversion: Converted from Radians to absolute Displacement (mm) and Phase Velocity (mm/s)."; p.level = 1

    # ---------------------------------------------------------
    # Slide 6: Heart Rate Extraction (Occlusion Filter)
    # ---------------------------------------------------------
    slide = prs.slides.add_slide(prs.slide_layouts[BULLET_SLIDE])
    slide.shapes.title.text = "Heart Rate Extraction Algorithm"
    tf = slide.placeholders[1].text_frame
    tf.text = "Overcoming the Cuff-Occlusion Problem:"
    p = tf.add_paragraph(); p.text = "Standard peak-counting fails because the artery is physically blocked (flatlined) during Phase 1."; p.level = 1
    p = tf.add_paragraph(); p.text = "Solution: Instantaneous Median Peak Tracking."; p.level = 1
    p = tf.add_paragraph(); p.text = "Algorithm dynamically measures the time interval between consecutive valid peaks."; p.level = 2
    p = tf.add_paragraph(); p.text = "Filters out massive temporal gaps (>1.5s) caused by occlusion."; p.level = 2
    p = tf.add_paragraph(); p.text = "Yields true instantaneous Heart Rate independent of the deflation schedule."; p.level = 1

    # ---------------------------------------------------------
    # Slide 7: Statistical Confirmation of Snaps
    # ---------------------------------------------------------
    slide = prs.slides.add_slide(prs.slide_layouts[BULLET_SLIDE])
    slide.shapes.title.text = "Statistical Domain: Kurtosis & TKEO"
    tf = slide.placeholders[1].text_frame
    tf.text = "How do we mathematically prove it's not random noise?"
    p = tf.add_paragraph(); p.text = "Sliding Kurtosis:"; p.level = 1
    p = tf.add_paragraph(); p.text = "Random noise is Gaussian (Kurtosis ~ 3)."; p.level = 2
    p = tf.add_paragraph(); p.text = "Arterial snaps are highly impulsive. Kurtosis spikes massively during the valid window."; p.level = 2
    p = tf.add_paragraph(); p.text = "Teager-Kaiser Energy Operator (TKEO):"; p.level = 1
    p = tf.add_paragraph(); p.text = "Measures true physical energy (Amplitude^2 * Frequency^2)."; p.level = 2
    p = tf.add_paragraph(); p.text = "Ignores slow movements (breathing) and violently highlights fast Korotkoff snaps."; p.level = 2

    # ---------------------------------------------------------
    # Slide 8: Spectrogram Analysis (Time-Frequency Domain)
    # ---------------------------------------------------------
    slide = prs.slides.add_slide(prs.slide_layouts[BULLET_SLIDE])
    slide.shapes.title.text = "Time-Frequency Domain (Spectrograms)"
    tf = slide.placeholders[1].text_frame
    tf.text = "Visualizing the frequency spread of mechanical impulses:"
    p = tf.add_paragraph(); p.text = "Short-Time Fourier Transform (STFT) with nperseg=4096 for ultra-high frequency resolution."; p.level = 1
    p = tf.add_paragraph(); p.text = "Wideband Analysis (0-150 Hz) proves physical wideband energy dispersion during snaps."; p.level = 1
    p = tf.add_paragraph(); p.text = "Eliminated derivative-amplification noise by running TFD on Displacement instead of Velocity."; p.level = 1
    p = tf.add_paragraph(); p.text = "Clearly shows Korotkoff harmonics fading out without artificial high-frequency electrical noise."; p.level = 1

    # ---------------------------------------------------------
    # Slide 9: Spectral Amplitude vs. Frequency
    # ---------------------------------------------------------
    slide = prs.slides.add_slide(prs.slide_layouts[BULLET_SLIDE])
    slide.shapes.title.text = "Amplitude Spectrum (FFT)"
    tf = slide.placeholders[1].text_frame
    tf.text = "Isolating the Active Korotkoff Window:"
    p = tf.add_paragraph(); p.text = "We extract data purely from the bounded SYS-DIA window and compute the FFT."; p.level = 1
    p = tf.add_paragraph(); p.text = "Magnitude Spectrum (0-60 Hz): Dominated entirely by massive spikes at the ~1-3 Hz heartbeat range."; p.level = 1
    p = tf.add_paragraph(); p.text = "Velocity Spectrum (0-60 Hz): Energy is physically shifted to the 20-50 Hz Korotkoff band."; p.level = 1
    p = tf.add_paragraph(); p.text = "Cross-validates that we are tracking genuine higher-frequency arterial resonance."; p.level = 1

    # ---------------------------------------------------------
    # Slide 10: Korotkoff Duration via Spectral Bandpower
    # ---------------------------------------------------------
    slide = prs.slides.add_slide(prs.slide_layouts[BULLET_SLIDE])
    slide.shapes.title.text = "Spectral Bandpower vs Time"
    tf = slide.placeholders[1].text_frame
    tf.text = "Defining strict duration directly from the frequency domain:"
    p = tf.add_paragraph(); p.text = "Integrated the STFT amplitude strictly across the 10-50 Hz target band."; p.level = 1
    p = tf.add_paragraph(); p.text = "Plotted as a continuous 1D amplitude curve over the entire 35-second recording."; p.level = 1
    p = tf.add_paragraph(); p.text = "Result: The 10-50Hz amplitude remains flat during occlusion, arcs perfectly through the active window, and drops at DIA."; p.level = 1
    p = tf.add_paragraph(); p.text = "Conversely, the 0.8-3.0 Hz (Heartbeat) Bandpower stays strong even after DIA, proving frequency separation."; p.level = 1

    # ---------------------------------------------------------
    # Slide 11: Physiological Synchronization
    # ---------------------------------------------------------
    slide = prs.slides.add_slide(prs.slide_layouts[BULLET_SLIDE])
    slide.shapes.title.text = "Physiological Synchronization Overlay"
    tf = slide.placeholders[1].text_frame
    tf.text = "The Ultimate Proof of Physiological Origin:"
    p = tf.add_paragraph(); p.text = "We zoomed in strictly on the 5-10 second Korotkoff active window."; p.level = 1
    p = tf.add_paragraph(); p.text = "Overlaid the low-frequency Heartbeat waveform (0.8-3.0 Hz) with the high-frequency Velocity Snaps (10-50 Hz)."; p.level = 1
    p = tf.add_paragraph(); p.text = "Observation: Every single high-frequency micro-vibration occurs precisely during the rising systolic edge of the heartbeat wave."; p.level = 1
    p = tf.add_paragraph(); p.text = "Conclusion: It is physically impossible for this to be random noise. It is 100% driven by the heart."; p.level = 1

    # ---------------------------------------------------------
    # Slide 12: The 22-Panel Dashboard
    # ---------------------------------------------------------
    slide = prs.slides.add_slide(prs.slide_layouts[BULLET_SLIDE])
    slide.shapes.title.text = "The 22-Panel Diagnostic Dashboard"
    tf = slide.placeholders[1].text_frame
    tf.text = "Comprehensive Visual Validation:"
    p = tf.add_paragraph(); p.text = "All of the aforementioned algorithms are compiled into a massive, automated 22-panel diagnostic dashboard."; p.level = 1
    p = tf.add_paragraph(); p.text = "Provides irrefutable multi-domain (Time, Statistical, Frequency, TFD) evidence."; p.level = 1
    p = tf.add_paragraph(); p.text = "Serves as the core analytical engine for processing future RMG clinical trials."; p.level = 1

    # ---------------------------------------------------------
    # Slide 13: Dashboard Image
    # ---------------------------------------------------------
    slide = prs.slides.add_slide(prs.slide_layouts[PIC_SLIDE])
    slide.shapes.title.text = "Results: Advanced 22-Panel RMG Dashboard"
    
    img_path = r'd:\Bioview\My_RF_work_v1\data_new\advanced_koro_validation_pressure_mapping_may12_2.png'
    if os.path.exists(img_path):
        # Add image, center it and scale it
        left = Inches(0.5)
        top = Inches(1.5)
        height = Inches(5.5)
        slide.shapes.add_picture(img_path, left, top, height=height)
    else:
        # Add text if image is missing
        txBox = slide.shapes.add_textbox(Inches(2), Inches(3), Inches(6), Inches(1))
        txBox.text_frame.text = "[Image not found. Place PNG here]"

    # ---------------------------------------------------------
    # Slide 14: Conclusion
    # ---------------------------------------------------------
    slide = prs.slides.add_slide(prs.slide_layouts[BULLET_SLIDE])
    slide.shapes.title.text = "Conclusion & Clinical Significance"
    tf = slide.placeholders[1].text_frame
    tf.text = "Summary of achievements:"
    p = tf.add_paragraph(); p.text = "Successfully implemented a non-invasive RF sensing pipeline for blood pressure tracking."; p.level = 1
    p = tf.add_paragraph(); p.text = "Eliminated phase accumulation errors and converted raw RF to absolute physical units (mm)."; p.level = 1
    p = tf.add_paragraph(); p.text = "Developed a robust bounding algorithm (SYS/DIA) verified by TKEO energy and Spectral Bandpower."; p.level = 1
    p = tf.add_paragraph(); p.text = "Provides higher fidelity and mathematical certainty compared to generic oscillometric cuffs."; p.level = 1

    # ---------------------------------------------------------
    # Slide 15: Future Work
    # ---------------------------------------------------------
    slide = prs.slides.add_slide(prs.slide_layouts[BULLET_SLIDE])
    slide.shapes.title.text = "Future Work"
    tf = slide.placeholders[1].text_frame
    tf.text = "Next steps for the RMG Pipeline:"
    p = tf.add_paragraph(); p.text = "Clinical Correlation: Map the extracted SYS/DIA indices to actual mmHg pressure readings from a reference gauge."; p.level = 1
    p = tf.add_paragraph(); p.text = "Real-Time Processing: Translate the offline Python dashboard into a real-time C++ / GNU Radio block."; p.level = 1
    p = tf.add_paragraph(); p.text = "Machine Learning: Train an inference model on the 10-50 Hz spectrogram features to automatically estimate absolute blood pressure without any cuff."; p.level = 1

    # Save presentation
    output_path = r'd:\Bioview\My_RF_work_v1\data_new\RMG_Korotkoff_Validation_Presentation.pptx'
    prs.save(output_path)
    print(f"Presentation generated successfully at: {output_path}")

if __name__ == '__main__':
    create_presentation()
